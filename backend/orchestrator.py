import os
from pathlib import Path
from typing import Optional
from dotenv import load_dotenv
from google import genai
from pptx import Presentation
from prompts.parser import parse_command
from file_searcher.search import find_files
from automation.pdfpipeline import PDFPipeline
from ppt_gen import generate_slides_content, create_pptx
from config import get_genai_model
from services import summarize as llm_summarize


load_dotenv(Path(__file__).parent.parent / ".env")

class AgentOrchestrator:
    def __init__(self):
        self.client = None
        self.model_name = get_genai_model()
        self.pipeline = PDFPipeline(summarizer=llm_summarize)

    def _get_client(self):
        if self.client is not None:
            return self.client

        api_key = os.environ.get("GEMINI_API_KEY")
        if not api_key:
            raise RuntimeError("GEMINI_API_KEY is not set")

        self.client = genai.Client(api_key=api_key)
        return self.client

    def run(self, query: str, state: Optional[dict] = None) -> dict:
        """Runs the orchestrator loop over a user query and returns a structured response."""
        if state is None:
            state = {}

        steps = []
        result_text = ""
        action_metadata = {}

        steps.append("Analyzing user query to understand intent...")
        parsed_action = parse_command(query)
        action = parsed_action.get("action", "chat").lower()
        file_name = parsed_action.get("file", "")
        
        steps.append(f"Detected action: '{action}' with target file: '{file_name or 'N/A'}'")
        action_metadata["parsed_action"] = parsed_action

        # ── GENERATE PPT ──────────────────────────────────────────────────────
        if action == "generate_ppt":
            topic = parsed_action.get("topic", "").strip()
            steps.append(f"Generating presentation{' on: ' + topic if topic else ''}...")

            source_text = None

            # Priority 1: a file was named → search and extract text
            if file_name:
                steps.append(f"Searching for '{file_name}' to build slides from...")
                matches = find_files(file_name)
                if matches:
                    target_path = matches[0]
                    steps.append(f"Found file: {target_path}")
                    try:
                        pipe_result = self.pipeline.process_file(target_path, summarize=False)
                        source_text = pipe_result.text
                        if not topic:
                            topic = os.path.splitext(pipe_result.filename)[0]
                    except Exception as e:
                        steps.append(f"Could not extract file text: {e}")

            # Priority 2: topic provided → ask Gemini to write content from scratch
            if not source_text and topic:
                steps.append(f"Generating slide content about '{topic}' using AI...")
                topic_prompt = f"""Write detailed content for a PowerPoint presentation about: "{topic}".
The presentation should have 6-8 slides.
Return ONLY a JSON array where each element has:
- "title": slide title (string)
- "bullets": list of 3-5 bullet points (list of strings)
Do not include markdown. Return raw JSON only."""
                try:
                    response = self._get_client().models.generate_content(
                        model=self.model_name,
                        contents=topic_prompt,
                    )
                    import json, re as _re
                    raw = (response.text or "").strip()
                    # Strip code fences if present
                    if raw.startswith("```"):
                        raw = _re.sub(r"^```[a-z]*\n?", "", raw)
                        raw = _re.sub(r"\n?```$", "", raw).strip()
                    slides_data = json.loads(raw)
                    ppt_filename = f"{topic.replace(' ', '_')[:40]}.pptx"
                    create_pptx(slides_data, ppt_filename)
                    state["ppt_filename"] = ppt_filename
                    steps.append(f"Presentation generated: {ppt_filename}")
                    result_text = f"✅ Your presentation on **{topic}** is ready!\n\nIt contains {len(slides_data)} slides. Click **Download PPT** to save it."
                except Exception as e:
                    steps.append(f"Error generating presentation: {e}")
                    result_text = f"Failed to generate presentation: {str(e)}"

            # Priority 3: no topic, no file → fall back to loaded session document
            elif not source_text and not topic:
                if state.get("text"):
                    source_text = state["text"]
                    topic = state.get("filename", "document")
                    steps.append("Using loaded document from session...")
                else:
                    result_text = "Please tell me what topic the presentation should be about. For example: **generate ppt on artificial intelligence**"
                    steps.append("No topic or document provided.")

            # If we have source_text (from file or session), generate slides from it
            if source_text and not state.get("ppt_filename"):
                try:
                    slides_data = generate_slides_content(source_text)
                    ppt_filename = f"{(topic or 'presentation').replace(' ', '_')[:40]}.pptx"
                    create_pptx(slides_data, ppt_filename)
                    state["ppt_filename"] = ppt_filename
                    steps.append(f"Presentation generated: {ppt_filename}")
                    result_text = f"✅ Your presentation on **{topic}** is ready!\n\nIt contains {len(slides_data)} slides. Click **Download PPT** to save it."
                except Exception as e:
                    steps.append(f"Error building PPT from document: {e}")
                    result_text = f"Failed to generate presentation: {str(e)}"

        elif action == "search" or action == "find":
            if not file_name:
                result_text = "Please specify a file name to search for."
                steps.append("Search failed: No file name provided.")
            else:
                steps.append(f"Searching for file matching '{file_name}'...")
                matches = find_files(file_name)
                if matches:
                    steps.append(f"Found {len(matches)} matching file(s).")
                    result_text = "Here are the files I found:\n" + "\n".join(f"- {m}" for m in matches)
                    
                    state["filename"] = os.path.basename(matches[0])
                    state["filepath"] = matches[0]
                else:
                    steps.append("No matching files found.")
                    result_text = f"I could not find any files matching '{file_name}'."

        elif action == "summarize" or action == "summary" or action == "summarise" or action == "brief" or action == "abstract" or action == " in short":
            target_path = None
            
            # If a filename was provided, search for it first
            if file_name:
                steps.append(f"Searching for file '{file_name}' to summarize...")
                matches = find_files(file_name)
                if matches:
                    target_path = matches[0]
                    steps.append(f"Found file at: {target_path}")
                else:
                    steps.append(f"Could not find '{file_name}' on disk. Checking current session...")

            # Fallback to loaded session file if no file search matches
            if not target_path and state.get("filepath"):
                target_path = state["filepath"]
                steps.append(f"Using loaded session file: {target_path}")

            if target_path:
                try:
                    steps.append(f"Processing PDF pipeline on '{os.path.basename(target_path)}'...")
                    try:
                        pipe_result = self.pipeline.process_file(target_path, summarize=True)
                    except Exception as e:
                        # If the file is not a PDF, try extracting text from PPTX as a fallback
                        if str(target_path).lower().endswith(('.pptx', '.ppt')):
                            try:

                                prs = Presentation(target_path)
                                slides_text = []
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if hasattr(shape, 'text') and shape.text:
                                            slides_text.append(shape.text)
                                text = "\n".join(slides_text)
                                # Build a minimal PipelineResult-like object
                                class _PR:
                                    def __init__(self, filename, text):
                                        self.filename = filename
                                        self.text = text
                                        self.summary = None
                                        self.page_count = 0
                                pipe_result = _PR(os.path.basename(target_path), text)
                                steps.append("Extracted text from PPTX fallback.")
                            except Exception as ppt_exc:
                                raise
                        else:
                            raise
                    
                    # Update active session state
                    state["text"] = pipe_result.text
                    state["filename"] = pipe_result.filename
                    state["filepath"] = target_path

                    steps.append("PDF text successfully extracted and summarized.")
                    result_text = f"### Summary of {pipe_result.filename}\n\n{pipe_result.summary}"
                except Exception as e:
                    steps.append(f"Error running pipeline: {str(e)}")
                    result_text = f"Failed to summarize file: {str(e)}"
            else:
                if state.get("text"):
                    steps.append("Summarizing current active document in memory...")
                    try:
                        from backend.services import summarize
                        summary = summarize(state["text"][:12000])
                        result_text = f"### Summary of loaded document\n\n{summary}"
                    except Exception as e:
                        result_text = f"Failed to summarize loaded document: {str(e)}"
                else:
                    steps.append("No active file or document loaded to summarize.")
                    result_text = "I couldn't find a file to summarize. Please load a PDF or specify its name in your message."

        elif action == "email" or action == "draft email" or action == "compose email":
            steps.append("Preparing to draft email...")
            doc_context = state.get("text", "")[:2000]
            
            prompt = f"""
            Draft a professional email based on the following query: "{query}".
            
            DOCUMENT CONTEXT:
            ---
            {doc_context}
            ---
            
            Format:
            [SUBJECT]
            <subject here>
            [BODY]
            <body here>
            """
            
            try:
                response = self._get_client().models.generate_content(
                    model=self.model_name,
                    contents=prompt
                )
                steps.append("Email draft generated successfully.")
                result_text = (response.text or "").strip()
            except Exception as e:
                steps.append(f"Error generating email: {str(e)}")
                result_text = f"Failed to generate email: {str(e)}"

        else:  
            doc_context = state.get("text", "")
            steps.append("Processing conversational query...")
            
            if doc_context:
                steps.append("Answering based on loaded document context...")
                prompt = f"""
                You are a professional workspace assistant. 
                Answer the user's question using ONLY the provided document text. If the answer is not in the text, say "I cannot find this in the document."
                
                USER QUESTION: {query}
                
                DOCUMENT TEXT:
                ---
                {doc_context[:6000]}
                ---
                """
            else:
                steps.append("No loaded document. Answering as general assistant...")
                prompt = query

            try:
                response = self._get_client().models.generate_content(
                    model=self.model_name,
                    contents=prompt
                )
                result_text = (response.text or "").strip()
                steps.append("Response generated.")
            except Exception as e:
                steps.append(f"Error communicating with AI: {str(e)}")
                result_text = f"AI request failed: {str(e)}"

        return {
            "result": result_text,
            "steps": steps,
            "state": state,
            "action_metadata": action_metadata,
            "ppt_filename": state.get("ppt_filename")
        }
