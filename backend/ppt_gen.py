from dotenv import load_dotenv
import os
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from google import genai
from google.genai import types
from config import get_genai_model

load_dotenv()
client = None


def _get_client():
    global client
    if client is not None:
        return client

    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY is not set")

    client = genai.Client(api_key=api_key)
    return client

# ─────────────────────────────────────────────
# Colour palette  (dark bg, light accent, text)
# Each tuple: (slide_bg, accent_bar, title_text, bullet_text, icon_color)
# ─────────────────────────────────────────────
THEMES = [
    # Deep navy + gold
    (RGBColor(0x0D, 0x1B, 0x2A), RGBColor(0xF4, 0xC4, 0x30), RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xE0, 0xE8, 0xF0), RGBColor(0xF4, 0xC4, 0x30)),
    # Dark forest + emerald
    (RGBColor(0x0F, 0x2D, 0x27), RGBColor(0x2E, 0xCC, 0x71), RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xD5, 0xF5, 0xE3), RGBColor(0x2E, 0xCC, 0x71)),
    # Deep plum + coral
    (RGBColor(0x1F, 0x0B, 0x2E), RGBColor(0xFF, 0x6B, 0x6B), RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xFA, 0xD7, 0xD7), RGBColor(0xFF, 0x6B, 0x6B)),
    # Midnight blue + cyan
    (RGBColor(0x04, 0x19, 0x35), RGBColor(0x00, 0xC8, 0xFF), RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xCC, 0xF0, 0xFF), RGBColor(0x00, 0xC8, 0xFF)),
    # Charcoal + amber
    (RGBColor(0x1C, 0x1C, 0x1E), RGBColor(0xFF, 0x9F, 0x0A), RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xFF, 0xED, 0xCC), RGBColor(0xFF, 0x9F, 0x0A)),
    # Dark teal + pink
    (RGBColor(0x07, 0x2A, 0x38), RGBColor(0xFF, 0x2D, 0x78), RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xFF, 0xCC, 0xE0), RGBColor(0xFF, 0x2D, 0x78)),
]

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

BULLET_ICONS = ["◆", "▸", "✦", "●", "➤", "★"]


# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────

def _fill_slide_bg(slide, color: RGBColor):
    """Set a solid background colour on the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor, transparency: int = 0):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()          # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.fore_color.theme_color  # keep rgb
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size, bold,
                 color: RGBColor, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


# ─────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────

def _build_title_slide(prs: Presentation, title: str, subtitle: str, theme: tuple):
    bg, accent, title_clr, bullet_clr, icon_clr = theme
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _fill_slide_bg(slide, bg)

    # Full-width accent bar at bottom
    _add_rect(slide, 0, Inches(6.6), SLIDE_W, Inches(0.9), accent)

    # Left decorative thick bar
    _add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, accent)

    # Right thick bar (thin)
    _add_rect(slide, SLIDE_W - Inches(0.08), 0, Inches(0.08), SLIDE_H,
              RGBColor(accent[0] // 2, accent[1] // 2, accent[2] // 2))

    # Main title
    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(1.8),
                 width=Inches(12.0), height=Inches(1.8),
                 text=title,
                 font_size=44, bold=True, color=title_clr,
                 align=PP_ALIGN.LEFT)

    # Accent line under title
    _add_rect(slide, Inches(0.6), Inches(3.7), Inches(4.5), Pt(4), accent)

    # Subtitle
    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(3.9),
                 width=Inches(11.5), height=Inches(1.2),
                 text=subtitle,
                 font_size=20, bold=False, color=bullet_clr,
                 align=PP_ALIGN.LEFT)

    # Bottom bar label
    _add_textbox(slide,
                 left=Inches(0.4), top=Inches(6.62),
                 width=Inches(6), height=Inches(0.5),
                 text="Powered by DocPilot AI",
                 font_size=13, bold=False,
                 color=RGBColor(0x1C, 0x1C, 0x1E),
                 align=PP_ALIGN.LEFT)


def _build_content_slide(prs: Presentation, title: str, bullets: list,
                         slide_num: int, theme: tuple):
    bg, accent, title_clr, bullet_clr, icon_clr = theme
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill_slide_bg(slide, bg)

    # ── Header bar ──
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.35), accent)

    # Slide number circle hint (top-right)
    num_box = slide.shapes.add_shape(9, SLIDE_W - Inches(1.0), Inches(0.15),
                                     Inches(0.75), Inches(0.75))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = bg
    num_box.line.fill.background()
    tf_num = num_box.text_frame
    tf_num.paragraphs[0].alignment = PP_ALIGN.CENTER
    run_num = tf_num.paragraphs[0].add_run()
    run_num.text = str(slide_num)
    run_num.font.size = Pt(14)
    run_num.font.bold = True
    run_num.font.color.rgb = accent
    run_num.font.name = "Calibri"

    # Title text on the bar
    _add_textbox(slide,
                 left=Inches(0.35), top=Inches(0.15),
                 width=Inches(11.8), height=Inches(1.0),
                 text=title,
                 font_size=30, bold=True, color=bg,
                 align=PP_ALIGN.LEFT)

    # ── Content area ──
    content_top = Inches(1.55)
    row_height = Inches(0.72)

    for i, bullet in enumerate(bullets[:6]):   # max 6 bullets
        y = content_top + i * row_height

        # Icon badge
        icon_shape = slide.shapes.add_shape(
            9,  # oval / ellipse-ish — we use rectangle with rounded feel
            Inches(0.28), y + Pt(3), Inches(0.38), Inches(0.38)
        )
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = accent
        icon_shape.line.fill.background()
        tf_icon = icon_shape.text_frame
        tf_icon.paragraphs[0].alignment = PP_ALIGN.CENTER
        run_icon = tf_icon.paragraphs[0].add_run()
        run_icon.text = BULLET_ICONS[i % len(BULLET_ICONS)]
        run_icon.font.size = Pt(11)
        run_icon.font.bold = True
        run_icon.font.color.rgb = bg
        run_icon.font.name = "Calibri"

        # Bullet text
        _add_textbox(slide,
                     left=Inches(0.8), top=y,
                     width=Inches(12.0), height=row_height,
                     text=bullet,
                     font_size=18, bold=False, color=bullet_clr,
                     align=PP_ALIGN.LEFT)

    # ── Footer line ──
    _add_rect(slide, 0, Inches(7.2), SLIDE_W, Pt(6), accent)


def _build_closing_slide(prs: Presentation, theme: tuple):
    bg, accent, title_clr, bullet_clr, icon_clr = theme
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, bg)

    # Diagonal accent block (fake — large rect bottom-right)
    _add_rect(slide, Inches(7.5), Inches(4.0), Inches(5.83), Inches(3.5), accent)

    # Left thick bar
    _add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, accent)

    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(2.2),
                 width=Inches(9), height=Inches(1.5),
                 text="Thank You",
                 font_size=52, bold=True, color=title_clr,
                 align=PP_ALIGN.LEFT)

    _add_rect(slide, Inches(0.6), Inches(3.85), Inches(3.5), Pt(5), accent)

    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(4.1),
                 width=Inches(9), height=Inches(0.9),
                 text="Generated by DocPilot AI",
                 font_size=18, bold=False, color=bullet_clr,
                 align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────

def generate_presentation(document_text: str) -> list:
    """Uses Gemini to structure key points into slide data."""
    if not document_text or not document_text.strip():
        return [{"title": "No Content Found",
                 "bullets": ["The uploaded file had no readable text.",
                             "Please ensure the PDF contains real text (not scanned images)."]}]

    prompt = f"""Based on the following document text, create a structured presentation of 5-50 slides.
    Return the response as a JSON array of slide objects.
    Each slide object MUST have:
    - "title": Title of the slide (string)
    - "bullets": An array of key bullet points (list of strings, 3-5 points per slide, each under 120 chars)
    Do not include markdown decoration. Return ONLY raw JSON.
    DOCUMENT TEXT:
    ---
    {document_text[:8000]}
    ---
    """
    try:
        model_name = get_genai_model()
        response = _get_client().models.generate_content(
            model=model_name,
            contents=prompt,
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )
        raw_text = (response.text or "").strip()

        if raw_text.startswith("```"):
            raw_text = re.sub(r"^```[a-z]*\n?", "", raw_text)
            raw_text = re.sub(r"\n?```$", "", raw_text).strip()

        try:
            return json.loads(raw_text)
        except json.JSONDecodeError:
            pass

        match = re.search(r"\[.*\]", raw_text, re.DOTALL)
        if match:
            return json.loads(match.group(0))

        raise ValueError(f"Could not parse JSON: {raw_text[:200]}")

    except Exception as e:
        print(f"Error generating slide content: {e}")
        return [{"title": "Document Summary",
                 "bullets": ["Automatic slide generation encountered an issue.",
                             "The document may contain scanned images instead of text.",
                             "Try uploading a PDF with selectable/copyable text."]}]


def generate_slides_content(document_text: str) -> list:
    return generate_presentation(document_text)


def create_pptx(slides_data: list, output_filename: str) -> str:
    """Build a visually designed PowerPoint and save to backend/static/."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Derive a document title from the first slide's title
    doc_title = slides_data[0].get("title", "Presentation") if slides_data else "Presentation"

    # Pick a theme palette (cycle through if many slides)
    theme = THEMES[0]

    # ── Title slide ──
    _build_title_slide(prs, doc_title, "AI-generated summary - DocPilot", theme)

    # ── Content slides ──
    for idx, slide_info in enumerate(slides_data):
        slide_theme = THEMES[idx % len(THEMES)]
        _build_content_slide(
            prs,
            title=slide_info.get("title", "Untitled"),
            bullets=slide_info.get("bullets", []),
            slide_num=idx + 1,
            theme=slide_theme,
        )

    # ── Closing slide ──
    _build_closing_slide(prs, THEMES[(len(slides_data)) % len(THEMES)])

    # Save
    static_dir = os.path.join(os.path.dirname(__file__), "static")
    os.makedirs(static_dir, exist_ok=True)
    filepath = os.path.join(static_dir, output_filename)
    prs.save(filepath)
    return filepath
